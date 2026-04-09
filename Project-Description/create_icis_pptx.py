"""
ICIS 2026 DSR Paper — 3-Slide Presentation
FIM Research Center corporate styling
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# === FIM Colour Palette ===
FIM_PRIMARY   = RGBColor(0xB1, 0x00, 0x34)  # Deep crimson
FIM_SECONDARY = RGBColor(0x4F, 0x78, 0x8C)  # Steel blue
FIM_TEXT      = RGBColor(0x64, 0x64, 0x64)   # Body text grey
FIM_DARK      = RGBColor(0x18, 0x18, 0x18)   # Near-black
FIM_LIGHT_BG  = RGBColor(0xF0, 0xF0, 0xF0)  # Light background
FIM_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
FIM_ACCENT    = RGBColor(0x9B, 0x9B, 0x9B)   # Muted accent

# Fonts (FIM brand: Exo for headings, Trebuchet MS for body)
FONT_HEADING = "Trebuchet MS"  # web-safe; Exo not guaranteed on all systems
FONT_BODY    = "Trebuchet MS"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)


# ──────────────────────────────────────────────────────────
# Helper functions
# ──────────────────────────────────────────────────────────

def add_background(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=FIM_TEXT, bold=False, alignment=PP_ALIGN.LEFT,
                font_name=FONT_BODY, line_spacing=1.3):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.line_spacing = Pt(int(font_size * line_spacing))
    return tf


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=FIM_TEXT, bold_first=False):
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = FONT_BODY
        p.level = 0
        p.space_before = Pt(6)
        p.line_spacing = Pt(int(font_size * 1.4))
        if bold_first and i == 0:
            p.font.bold = True
    return tf


def add_accent_bar(slide, left, top, width, height, color=FIM_PRIMARY):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_card(slide, left, top, width, height, title, body_items,
             accent_color=FIM_SECONDARY):
    # Card background
    card = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    card.fill.solid()
    card.fill.fore_color.rgb = FIM_WHITE
    card.line.color.rgb = FIM_LIGHT_BG
    card.line.width = Pt(1)
    card.shadow.inherit = False

    # Accent top bar
    add_accent_bar(slide, left, top, width, 0.06, accent_color)

    # Title
    add_textbox(slide, left + 0.25, top + 0.2, width - 0.5, 0.5,
                title, font_size=16, color=accent_color, bold=True)

    # Body bullets
    add_bullet_list(slide, left + 0.25, top + 0.7, width - 0.5, height - 1.0,
                    body_items, font_size=13, color=FIM_TEXT)


def add_footer(slide, text="FIM Research Center | ICIS 2026"):
    add_textbox(slide, 0.5, 7.0, 12.3, 0.4, text,
                font_size=10, color=FIM_ACCENT, alignment=PP_ALIGN.RIGHT)


def add_slide_number(slide, num):
    add_textbox(slide, 12.5, 7.0, 0.5, 0.4, str(num),
                font_size=10, color=FIM_ACCENT, alignment=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════
# SLIDE 1 — Title & Research Overview
# ══════════════════════════════════════════════════════════
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_background(slide1, FIM_WHITE)

# Top accent bar
add_accent_bar(slide1, 0, 0, 13.333, 0.12, FIM_PRIMARY)

# Left crimson block for visual weight
add_accent_bar(slide1, 0, 0.12, 0.08, 7.38, FIM_PRIMARY)

# Title
add_textbox(slide1, 0.8, 0.4, 11.5, 1.2,
            "All Carriers Allowed",
            font_size=36, color=FIM_PRIMARY, bold=True,
            font_name=FONT_HEADING, alignment=PP_ALIGN.LEFT)

# Subtitle
add_textbox(slide1, 0.8, 1.2, 11.5, 0.8,
            "Design Principles for DLT-Based Cross-Domain Data Sharing\nBetween Guarantee of Origin Schemes",
            font_size=20, color=FIM_SECONDARY, bold=False,
            font_name=FONT_HEADING)

# Divider line
add_accent_bar(slide1, 0.8, 2.15, 5.0, 0.03, FIM_SECONDARY)

# Authors & venue
add_textbox(slide1, 0.8, 2.35, 6.0, 0.5,
            "Jonathan Albrecht  |  ICIS 2026  |  FIM Research Center",
            font_size=14, color=FIM_ACCENT)

# ── Research Question box ──
rq_bg = slide1.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(0.8), Inches(3.1), Inches(11.5), Inches(1.2)
)
rq_bg.fill.solid()
rq_bg.fill.fore_color.rgb = FIM_LIGHT_BG
rq_bg.line.fill.background()

add_textbox(slide1, 1.0, 3.15, 0.6, 0.4, "RQ",
            font_size=14, color=FIM_WHITE, bold=True,
            alignment=PP_ALIGN.CENTER)
# RQ label pill
rq_pill = slide1.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(1.0), Inches(3.15), Inches(0.55), Inches(0.35)
)
rq_pill.fill.solid()
rq_pill.fill.fore_color.rgb = FIM_PRIMARY
rq_pill.line.fill.background()
rq_pill_tf = rq_pill.text_frame
rq_pill_tf.paragraphs[0].text = "RQ"
rq_pill_tf.paragraphs[0].font.size = Pt(12)
rq_pill_tf.paragraphs[0].font.color.rgb = FIM_WHITE
rq_pill_tf.paragraphs[0].font.bold = True
rq_pill_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
rq_pill_tf.paragraphs[0].font.name = FONT_HEADING
rq_pill_tf.vertical_anchor = MSO_ANCHOR.MIDDLE

add_textbox(slide1, 1.7, 3.15, 10.3, 1.0,
            "How to design an information system of interoperable GO schemes that allows\n"
            "verifiable cross-domain data sharing and asset transfer to enable energy\n"
            "carrier conversion processes in multi-carrier energy systems?",
            font_size=16, color=FIM_DARK, bold=False)

# ── Three content columns ──
# Column 1: Paper Structure
add_card(slide1, 0.8, 4.6, 3.5, 2.2,
         "Paper Structure",
         [
             "Design Science Research (Hevner 2004)",
             "8 iterative design cycles",
             "Multi-vocal literature review",
             "4 design principles (Gregor et al. 2020)",
         ],
         accent_color=FIM_PRIMARY)

# Column 2: Research Contribution
add_card(slide1, 4.7, 4.6, 3.5, 2.2,
         "Contributions",
         [
             "4 empirically grounded design principles",
             "Full-stack HLF prototype (v8.0)",
             "Caliper benchmarks: 50 TPS write",
             "Nascent design theory for GO systems",
         ],
         accent_color=FIM_SECONDARY)

# Column 3: Practical Aim
add_card(slide1, 8.6, 4.6, 3.5, 2.2,
         "Aim & Scope",
         [
             "RED III mandates multi-carrier GOs",
             "Cross-carrier conversion issuance",
             "Privacy + verifiability trade-off",
             "Scalable to EU-wide deployment",
         ],
         accent_color=FIM_PRIMARY)

add_footer(slide1)
add_slide_number(slide1, 1)

# Speaker notes
slide1.notes_slide.notes_text_frame.text = (
    "This paper addresses the design challenge of interoperable Guarantee of Origin "
    "schemes across multiple energy carriers (electricity, hydrogen, biogas, heating). "
    "RED III mandates cross-carrier GO support. We use Design Science Research to derive "
    "4 design principles and instantiate them in a Hyperledger Fabric prototype developed "
    "over 8 iterative cycles. GitHub repository: to be linked in final version."
)


# ══════════════════════════════════════════════════════════
# SLIDE 2 — Hyperledger Fabric Deep Dive
# ══════════════════════════════════════════════════════════
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide2, FIM_WHITE)

# Top accent bar
add_accent_bar(slide2, 0, 0, 13.333, 0.12, FIM_SECONDARY)

# Left accent bar
add_accent_bar(slide2, 0, 0.12, 0.08, 7.38, FIM_SECONDARY)

# Section title
add_textbox(slide2, 0.8, 0.3, 11.5, 0.6,
            "Hyperledger Fabric Prototype — Why, How & What We Achieved",
            font_size=28, color=FIM_SECONDARY, bold=True,
            font_name=FONT_HEADING)

# Divider
add_accent_bar(slide2, 0.8, 1.0, 4.0, 0.03, FIM_PRIMARY)

# ── Left column: WHY ──
add_textbox(slide2, 0.8, 1.3, 3.7, 0.4,
            "Why Hyperledger Fabric?", font_size=18,
            color=FIM_PRIMARY, bold=True, font_name=FONT_HEADING)

add_bullet_list(slide2, 0.8, 1.8, 3.7, 2.2, [
    "Permissioned: known participants (national issuing bodies)",
    "Private data collections for confidential GO attributes",
    "Channel architecture for carrier/region data partitioning",
    "Execute-order-validate flow for deterministic smart contracts",
    "Raft consensus: crash-fault-tolerant, enterprise-grade",
], font_size=13, color=FIM_TEXT)

# ── Middle column: HOW ──
add_textbox(slide2, 5.0, 1.3, 3.7, 0.4,
            "How — 8 Iterative Design Cycles", font_size=18,
            color=FIM_PRIMARY, bold=True, font_name=FONT_HEADING)

# Version timeline cards
versions = [
    ("v1–v3", "Master thesis: conceptual design,\nbasic GO lifecycle, refinements"),
    ("v4–v5", "Scalability (pagination, tombstones),\nCEN-EN 16325 alignment, CQRS"),
    ("v6–v7", "Crypto-secure commitments,\nbridge protocol, smart meter ECDSA"),
    ("v8", "Multi-channel architecture:\ncarrier-per-channel, Lock→Mint→Finalize"),
]

y_pos = 1.85
for label, desc in versions:
    # Version pill
    pill = slide2.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.0), Inches(y_pos), Inches(0.65), Inches(0.3)
    )
    pill.fill.solid()
    pill.fill.fore_color.rgb = FIM_SECONDARY
    pill.line.fill.background()
    ptf = pill.text_frame
    ptf.paragraphs[0].text = label
    ptf.paragraphs[0].font.size = Pt(10)
    ptf.paragraphs[0].font.color.rgb = FIM_WHITE
    ptf.paragraphs[0].font.bold = True
    ptf.paragraphs[0].font.name = FONT_HEADING
    ptf.paragraphs[0].alignment = PP_ALIGN.CENTER
    ptf.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_textbox(slide2, 5.8, y_pos - 0.02, 2.9, 0.6,
                desc, font_size=11, color=FIM_TEXT)
    y_pos += 0.72

# ── Right column: WHAT WE ACHIEVED ──
add_textbox(slide2, 9.2, 1.3, 3.8, 0.4,
            "What We Achieved", font_size=18,
            color=FIM_PRIMARY, bold=True, font_name=FONT_HEADING)

# Achievement stat cards
stats = [
    ("50 TPS", "sustained write throughput"),
    ("2,000 TPS", "read throughput"),
    ("10", "smart contract namespaces"),
    ("~50", "exported chaincode functions"),
    ("4-org", "Raft network on Hetzner VM"),
]

y_stat = 1.85
for val, label in stats:
    # Stat value
    add_textbox(slide2, 9.2, y_stat, 1.5, 0.35,
                val, font_size=18, color=FIM_PRIMARY, bold=True,
                font_name=FONT_HEADING)
    # Stat label
    add_textbox(slide2, 10.6, y_stat + 0.03, 2.4, 0.35,
                label, font_size=12, color=FIM_TEXT)
    y_stat += 0.5

# ── Bottom section: Full-Stack Architecture summary ──
bottom_bg = slide2.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(0.8), Inches(4.8), Inches(11.7), Inches(2.0)
)
bottom_bg.fill.solid()
bottom_bg.fill.fore_color.rgb = FIM_LIGHT_BG
bottom_bg.line.fill.background()

add_textbox(slide2, 1.0, 4.9, 3.5, 0.4,
            "Full-Stack Architecture", font_size=16,
            color=FIM_SECONDARY, bold=True, font_name=FONT_HEADING)

# Architecture layers
arch_items = [
    ("Chaincode (Go)", "10 namespaces: device, go, conversion,\nbridge, market, plausibility, …"),
    ("Backend (Node.js/TS)", "REST API, Fabric Gateway SDK,\nWallet management, event listeners"),
    ("Frontend (React/TS)", "Dashboard, GO lifecycle mgmt,\ncross-carrier conversion UI"),
]

x_arch = 1.0
for title, desc in arch_items:
    add_textbox(slide2, x_arch, 5.35, 3.3, 0.3,
                title, font_size=13, color=FIM_PRIMARY, bold=True)
    add_textbox(slide2, x_arch, 5.65, 3.3, 0.8,
                desc, font_size=11, color=FIM_TEXT)
    x_arch += 3.7

# Scalability note
add_textbox(slide2, 1.0, 6.45, 11.0, 0.4,
            "Scalability tested on production-grade Hetzner VM (16 vCPU, 32 GB RAM) "
            "with Hyperledger Caliper v0.6.0 benchmarks  |  "
            "GitHub: github.com/[to-be-added]",
            font_size=11, color=FIM_ACCENT)

add_footer(slide2)
add_slide_number(slide2, 2)

slide2.notes_slide.notes_text_frame.text = (
    "The prototype was built over 8 iterative design cycles. v1-v3 during the master thesis "
    "(2 years ago), v4-v7 extending to production-grade quality, v8 introducing multi-channel "
    "architecture. The full-stack includes Go chaincode (10 namespaces, ~50 functions), "
    "Node.js/TypeScript backend, and React/TypeScript frontend. Benchmarked on a Hetzner VM "
    "using Hyperledger Caliper. GitHub link will be in the paper — blind review no longer required. "
    "Note: 8 semi-structured interviews were conducted during the master thesis phase to elicit "
    "requirements from domain experts, but these were not used for artifact evaluation — "
    "only for initial requirements gathering. Consider including them in the paper for completeness."
)


# ══════════════════════════════════════════════════════════
# SLIDE 3 — Questions / Discussion
# ══════════════════════════════════════════════════════════
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide3, FIM_WHITE)

# Full-width crimson block top section
top_block = slide3.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(0), Inches(0), Inches(13.333), Inches(3.2)
)
top_block.fill.solid()
top_block.fill.fore_color.rgb = FIM_PRIMARY
top_block.line.fill.background()

add_textbox(slide3, 0.8, 0.8, 11.5, 1.2,
            "Questions & Discussion",
            font_size=44, color=FIM_WHITE, bold=True,
            font_name=FONT_HEADING, alignment=PP_ALIGN.LEFT)

add_textbox(slide3, 0.8, 2.0, 11.5, 0.6,
            "Jonathan Albrecht  |  FIM Research Center  |  ICIS 2026",
            font_size=18, color=RGBColor(0xFF, 0xBC, 0x7D), bold=False)

# Discussion points in lower section
add_textbox(slide3, 0.8, 3.6, 11.5, 0.5,
            "Discussion Points",
            font_size=22, color=FIM_SECONDARY, bold=True,
            font_name=FONT_HEADING)

add_accent_bar(slide3, 0.8, 4.15, 3.0, 0.03, FIM_SECONDARY)

# Two-column discussion layout
disc_left = [
    "GitHub repository will be linked in paper\n(no longer under blind review)",
    "8 expert interviews from master thesis phase\n→ requirements elicitation, not artifact evaluation",
    "Consider including interviews for\nmethodological completeness",
]

disc_right = [
    "How generalizable are the design principles\nbeyond the EU GO market context?",
    "Channel-per-carrier vs. channel-per-region\npartitioning trade-offs?",
    "Next steps: pilot with national issuing body,\nreal-world metering device integration",
]

y_disc = 4.4
for item in disc_left:
    # Bullet dot
    dot = slide3.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(1.0), Inches(y_disc + 0.07), Inches(0.12), Inches(0.12)
    )
    dot.fill.solid()
    dot.fill.fore_color.rgb = FIM_PRIMARY
    dot.line.fill.background()
    add_textbox(slide3, 1.3, y_disc, 5.0, 0.7,
                item, font_size=14, color=FIM_TEXT)
    y_disc += 0.73

y_disc = 4.4
for item in disc_right:
    dot = slide3.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(7.0), Inches(y_disc + 0.07), Inches(0.12), Inches(0.12)
    )
    dot.fill.solid()
    dot.fill.fore_color.rgb = FIM_SECONDARY
    dot.line.fill.background()
    add_textbox(slide3, 7.3, y_disc, 5.5, 0.7,
                item, font_size=14, color=FIM_TEXT)
    y_disc += 0.73

# Contact info at bottom
add_accent_bar(slide3, 0.8, 6.7, 11.7, 0.02, FIM_LIGHT_BG)

add_textbox(slide3, 0.8, 6.8, 11.5, 0.3,
            "fim-rc.de  |  Universität Bayreuth & Technische Hochschule Augsburg",
            font_size=11, color=FIM_ACCENT, alignment=PP_ALIGN.CENTER)

add_footer(slide3, "")
add_slide_number(slide3, 3)

slide3.notes_slide.notes_text_frame.text = (
    "Key points to address:\n"
    "1) GitHub link: Paper is past blind review stage, so including the repo URL is fine.\n"
    "2) Interviews: 8 semi-structured interviews were conducted 2 years ago during the master "
    "thesis to gather requirements. They were NOT used for artifact evaluation — only for "
    "initial requirement elicitation. Including them adds methodological depth but they are "
    "not strictly necessary for the DSR contribution.\n"
    "3) Open for questions on design principles, Fabric architecture, or evaluation approach."
)


# ══════════════════════════════════════════════════════════
# Save
# ══════════════════════════════════════════════════════════
output_path = r"C:\Users\jona\Algo-Trading\Master-Thesis\HLF-GOconversionissuance-JA-MA\Project-Description\20260410_ICIS_2026_FIM_Presentation.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
